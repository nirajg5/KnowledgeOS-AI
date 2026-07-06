"""
Document Processor

Extracts plain text from supported documents.
"""

from pathlib import Path

import fitz
import pandas as pd
import markdown
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


class DocumentProcessor:

    @staticmethod
    def extract_pdf(file_path: Path):

        document = fitz.open(file_path)

        text = ""

        for page in document:
            text += page.get_text()

        document.close()

        return text


    @staticmethod
    def extract_docx(file_path: Path):

        document = Document(file_path)

        return "\n".join(
            paragraph.text
            for paragraph in document.paragraphs
        )


    @staticmethod
    def extract_txt(file_path: Path):

        return file_path.read_text(
            encoding="utf-8",
            errors="ignore"
        )


    @staticmethod
    def extract_csv(file_path: Path):

        dataframe = pd.read_csv(file_path)

        return dataframe.to_string(index=False)


    @staticmethod
    def extract_xlsx(file_path: Path):

        workbook = load_workbook(
            file_path,
            data_only=True
        )

        text = ""

        for sheet in workbook.worksheets:

            text += f"\nSheet : {sheet.title}\n"

            for row in sheet.iter_rows(values_only=True):

                values = [
                    str(value)
                    for value in row
                    if value is not None
                ]

                text += " ".join(values)

                text += "\n"

        return text


    @staticmethod
    def extract_pptx(file_path: Path):

        presentation = Presentation(file_path)

        text = ""

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

        return text


    @staticmethod
    def extract_markdown(file_path: Path):

        md = file_path.read_text(
            encoding="utf-8"
        )

        html = markdown.markdown(md)

        return BeautifulSoup(
            html,
            "html.parser"
        ).get_text()


    @staticmethod
    def extract_text(file_path: str):

        path = Path(file_path)

        extension = path.suffix.lower()

        if extension == ".pdf":

            return DocumentProcessor.extract_pdf(path)

        elif extension in [".docx", ".doc"]:

            return DocumentProcessor.extract_docx(path)

        elif extension == ".txt":

            return DocumentProcessor.extract_txt(path)

        elif extension == ".csv":

            return DocumentProcessor.extract_csv(path)

        elif extension in [".xls", ".xlsx"]:

            return DocumentProcessor.extract_xlsx(path)

        elif extension in [".ppt", ".pptx"]:

            return DocumentProcessor.extract_pptx(path)

        elif extension == ".md":

            return DocumentProcessor.extract_markdown(path)

        raise ValueError(
            f"Unsupported file type : {extension}"
        )